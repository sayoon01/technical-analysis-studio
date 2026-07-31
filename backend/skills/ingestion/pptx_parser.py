"""PPTX ingestion — each slide becomes one RawPage."""

from __future__ import annotations

from pathlib import Path

from backend.skills.ingestion.pdf_parser import RawPage, RawTextBlock
from backend.skills.ingestion.text_parser import TextDocument


def parse_pptx(path: str | Path) -> TextDocument:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        raise RuntimeError(
            "python-pptx is required for PPTX ingestion. pip install python-pptx"
        ) from e

    prs = Presentation(str(path))
    # EMU → approximate points (914400 EMU = 1 inch, 72 pt = 1 inch)
    width = float(prs.slide_width) / 914400.0 * 72.0 if prs.slide_width else 720.0
    height = float(prs.slide_height) / 914400.0 * 72.0 if prs.slide_height else 540.0

    pages: list[RawPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        blocks: list[RawTextBlock] = []
        texts: list[str] = []
        image_count = 0
        drawing_count = 0

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                paras = []
                for p in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in p.runs).strip()
                    if not line and p.text:
                        line = p.text.strip()
                    if line:
                        paras.append(line)
                text = "\n".join(paras).strip()
                if text:
                    texts.append(text)
                    blocks.append(
                        RawTextBlock(
                            text=text,
                            bbox=_shape_bbox(shape, width, height),
                            block_type="TEXT",
                            confidence=1.0,
                        )
                    )

            if getattr(shape, "has_table", False) and shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    rows.append(" | ".join(cells))
                table_text = "\n".join(rows).strip()
                if table_text:
                    texts.append(table_text)
                    blocks.append(
                        RawTextBlock(
                            text=table_text,
                            bbox=_shape_bbox(shape, width, height),
                            block_type="TABLE",
                            confidence=1.0,
                        )
                    )

            st = getattr(shape, "shape_type", None)
            if st == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            else:
                drawing_types = {
                    getattr(MSO_SHAPE_TYPE, name)
                    for name in (
                        "AUTO_SHAPE",
                        "FREEFORM",
                        "LINE",
                        "CONNECTOR",
                        "GROUP",
                    )
                    if hasattr(MSO_SHAPE_TYPE, name)
                }
                if st in drawing_types:
                    drawing_count += 1

        # notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[notes]\n{notes}")
                blocks.append(
                    RawTextBlock(
                        text=notes,
                        bbox=(0.0, height - 40.0, width, height),
                        block_type="TEXT",
                        confidence=0.9,
                    )
                )

        full = "\n\n".join(texts)
        if not blocks and not full:
            blocks.append(
                RawTextBlock(
                    text="",
                    bbox=(0.0, 0.0, width, height),
                    block_type="TEXT",
                    confidence=0.0,
                )
            )

        pages.append(
            RawPage(
                page_number=i,
                width=width,
                height=height,
                text_layer_available=bool(full.strip()),
                full_text=full,
                blocks=blocks,
                image_count=image_count,
                drawing_count=drawing_count,
            )
        )

    if not pages:
        pages.append(
            RawPage(
                page_number=1,
                width=width,
                height=height,
                text_layer_available=False,
                full_text="",
                blocks=[],
            )
        )
    return TextDocument(pages=pages)


def _shape_bbox(shape, page_w: float, page_h: float) -> tuple[float, float, float, float]:
    try:
        x0 = float(shape.left) / 914400.0 * 72.0
        y0 = float(shape.top) / 914400.0 * 72.0
        x1 = x0 + float(shape.width) / 914400.0 * 72.0
        y1 = y0 + float(shape.height) / 914400.0 * 72.0
        return (x0, y0, min(x1, page_w), min(y1, page_h))
    except Exception:
        return (0.0, 0.0, page_w, page_h)
