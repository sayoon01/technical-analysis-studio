"""Ingestion for PPTX / XLSX / CSV (2차 포맷)."""

from __future__ import annotations

import sqlite3
from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from backend.services.project_service import ProjectService, SourceService
from backend.skills.ingestion.file_detector import detect_format, is_supported
from backend.skills.ingestion.pptx_parser import parse_pptx
from backend.skills.ingestion.spreadsheet_parser import parse_spreadsheet
from backend.storage.database import init_schema
from backend.storage.file_store import FileStore


def _build_pptx(path: Path) -> Path:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    box = s1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    box.text_frame.text = "NPU 추론 가속 아키텍처 개요"

    s2 = prs.slides.add_slide(blank)
    box2 = s2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8), Inches(0.6))
    box2.text_frame.text = "정량 성과"
    table = s2.shapes.add_table(3, 3, Inches(0.5), Inches(1.2), Inches(8), Inches(2)).table
    table.cell(0, 0).text = "지표"
    table.cell(0, 1).text = "변화"
    table.cell(0, 2).text = "비고"
    table.cell(1, 0).text = "처리량"
    table.cell(1, 1).text = "+32%"
    table.cell(1, 2).text = "batch=8"
    table.cell(2, 0).text = "지연시간"
    table.cell(2, 1).text = "-18%"
    table.cell(2, 2).text = "p95"

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def _build_xlsx(path: Path) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "metrics"
    ws.append(["metric", "change", "unit"])
    ws.append(["생산량", 8, "%"])
    ws.append(["클레임", -60, "%"])
    ws2 = wb.create_sheet("notes")
    ws2.append(["기존 MES 수기 입력 문제"])
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    return path


def test_detect_pptx_xlsx_csv():
    assert detect_format("a.pptx") == "pptx"
    assert detect_format("b.xlsx") == "xlsx"
    assert detect_format("c.csv") == "csv"
    assert is_supported("d.xlsm")
    assert detect_format("e.hwpx") == "hwpx"
    assert detect_format("f.hwp") == "hwp"


def test_parse_hwpx_text(tmp_path: Path):
    from backend.skills.ingestion.hwp_parser import parse_hwp
    import zipfile

    hwpx = tmp_path / "sample.hwpx"
    xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        "<section><p>클라우드 MES 도입 개요</p><p>처리량 +12%</p></section>"
    )
    with zipfile.ZipFile(hwpx, "w") as zf:
        zf.writestr("Contents/section0.xml", xml)
    doc = parse_hwp(hwpx)
    assert any("클라우드" in p.full_text for p in doc.pages)
    assert any("처리량" in p.full_text for p in doc.pages)


def test_parse_pptx_slides(tmp_path: Path):
    path = _build_pptx(tmp_path / "sample.pptx")
    doc = parse_pptx(path)
    assert len(doc.pages) == 2
    assert "NPU" in doc.pages[0].full_text
    assert "처리량" in doc.pages[1].full_text
    assert any(b.block_type == "TABLE" for b in doc.pages[1].blocks)


def test_parse_xlsx_sheets(tmp_path: Path):
    path = _build_xlsx(tmp_path / "sample.xlsx")
    doc = parse_spreadsheet(path)
    assert len(doc.pages) == 2
    assert "생산량" in doc.pages[0].full_text
    assert any(b.block_type == "TABLE" for b in doc.pages[0].blocks)
    assert "수기" in doc.pages[1].full_text


def test_pipeline_process_pptx_xlsx(tmp_path: Path):
    db = tmp_path / "t.db"
    data = tmp_path / "data"
    data.mkdir()
    conn = sqlite3.connect(str(db))
    conn.row_factory = sqlite3.Row
    init_schema(conn)
    store = FileStore(root=data / "projects")
    projects = ProjectService(conn)
    sources = SourceService(conn, store)

    project = projects.create("pptx-xlsx")
    pid = project["project_id"]

    pptx = _build_pptx(tmp_path / "deck.pptx")
    xlsx = _build_xlsx(tmp_path / "book.xlsx")

    u1 = sources.upload(pid, "deck.pptx", pptx.read_bytes())
    r1 = sources.process(u1["source_id"])
    assert r1["page_count"] == 2
    assert r1["status"] == "READY"
    page = sources.get_page(u1["source_id"], 1)
    assert "NPU" in (page.get("text") or "")

    u2 = sources.upload(pid, "book.xlsx", xlsx.read_bytes())
    r2 = sources.process(u2["source_id"])
    assert r2["page_count"] == 2
    assert r2["block_count"] >= 2
